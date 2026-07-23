"""Office 文档读取与生成服务，以及“卷宗”（问策 Run 最终报告）读取/列举/回填。

所有路径都限制在当前团队工作目录。卷宗对应一次 GtTaskRun 的最终综合报告：
持久化为团队 outputs 目录下的 Markdown 文件（``综合分析报告-run-{id}.md``），
其相对路径记录在 run.metadata['final_report_path']。历史 Run（在 H1 修复
_write_final_report_artifact 之前完成）可能只有 DB 中的 final_answer 而没有落盘
文件，导致“历史卷宗无法查看/下载”。这里的 get_dossier 在读取时按需回填该文件，
使历史卷宗恢复可查可下载。"""
from __future__ import annotations
import asyncio, csv, json, logging, os, re
from pathlib import Path
from dal.db import gtTaskRunManager, gtTeamManager
from model.dbModel.gtTaskRun import GtTaskRun
from service.roomService import ToolCallContext
from util import configUtil, fileUtil

logger = logging.getLogger(__name__)


def _safe_name(value:str)->str:
    return re.sub(r'[^\w\u4e00-\u9fff.-]+','_',value.strip())[:80] or '数字人生文档'

async def _team_workdir(context:ToolCallContext)->Path:
    return await _team_workdir_by_id(context.team_id)

async def _team_workdir_by_id(team_id:int)->Path:
    team=await gtTeamManager.get_team_by_id(team_id)
    if team is None: raise ValueError('当前团队不存在')
    root=configUtil.get_app_config().setting.workspace_root
    if not root: raise ValueError('团队工作目录尚未配置')
    workdir=Path((team.config or {}).get('working_directory') or os.path.join(root,team.name)).resolve()
    workdir.mkdir(parents=True,exist_ok=True)
    return workdir

def _resolve(workdir:Path,relative_path:str)->Path:
    path=(workdir/relative_path).resolve();fileUtil.assert_path_within_sandbox(str(path),str(workdir));return path

# 单文件解析大小上限：超出直接拒绝，避免大文件解析长时间占用线程池与内存
_EXTRACT_MAX_FILE_BYTES = 20 * 1024 * 1024


def _parse_office_file_text(source:Path, ext:str)->str | None:
    """同步解析 Office/文本文件正文；不支持的扩展名返回 None。

    CPU/IO 密集（docx/xlsx/pptx 解析、大文本读取），必须经 asyncio.to_thread 调用，
    不得在事件循环内直接执行（审计 M8）。
    """
    parts=[]
    if ext in {'.md','.markdown','.txt','.json','.csv'}:
        return source.read_text(encoding='utf-8',errors='replace')
    elif ext=='.docx':
        from docx import Document
        doc=Document(source);parts.extend(p.text for p in doc.paragraphs if p.text.strip())
        for table in doc.tables:
            parts.extend(' | '.join(cell.text for cell in row.cells) for row in table.rows)
        return '\n'.join(parts)
    elif ext in {'.xlsx','.xlsm'}:
        from openpyxl import load_workbook
        book=load_workbook(source,data_only=False,read_only=True)
        for sheet in book.worksheets:
            parts.append(f'## 工作表：{sheet.title}')
            for row in sheet.iter_rows(values_only=True):parts.append(' | '.join('' if v is None else str(v) for v in row))
        return '\n'.join(parts)
    elif ext=='.pptx':
        from pptx import Presentation
        deck=Presentation(source)
        for index,slide in enumerate(deck.slides,1):
            parts.append(f'## 第 {index} 页')
            parts.extend(shape.text for shape in slide.shapes if hasattr(shape,'text') and shape.text.strip())
        return '\n'.join(parts)
    return None

async def extract_office_file(path:str,_context:ToolCallContext=None)->dict:
    """读取已上传的 Word、Excel、PPT、Markdown 或文本文件并返回可供分析的正文。

    Args:
        path: 团队工作目录中的相对路径，如 uploads/20260711_report.docx。
    """
    if _context is None:return {'success':False,'message':'缺少团队上下文'}
    workdir=await _team_workdir(_context);source=_resolve(workdir,path)
    if not source.is_file():return {'success':False,'message':f'未找到文件: {path}'}
    size=source.stat().st_size
    if size>_EXTRACT_MAX_FILE_BYTES:return {'success':False,'message':f'文件过大（{size // 1024 // 1024}MB，上限 20MB），请先拆分或转换'}
    ext=source.suffix.lower()
    text=await asyncio.to_thread(_parse_office_file_text, source, ext)
    if text is None:return {'success':False,'message':f'当前不能直接解析 {ext}；请先转换为 docx/xlsx/pptx/md'}
    limit=120000;return {'success':True,'path':path,'characters':len(text),'truncated':len(text)>limit,'content':text[:limit]}

async def generate_office_file(format:str,title:str,content:str,filename:str='',_context:ToolCallContext=None)->dict:
    """生成 Word、Excel、PPT 或 Markdown 文件并保存到 outputs 目录。

    Args:
        format: docx、xlsx、pptx 或 md。
        title: 文档标题。
        content: Markdown 风格正文；生成 Excel 时每行可用竖线分隔列，生成 PPT 时用二级标题分隔页面。
        filename: 可选文件名，不含目录。
    """
    if _context is None:return {'success':False,'message':'缺少团队上下文'}
    fmt=format.lower().lstrip('.');allowed={'docx','xlsx','pptx','md'}
    if fmt not in allowed:return {'success':False,'message':'仅支持 docx、xlsx、pptx、md'}
    workdir=await _team_workdir(_context);out=workdir/'outputs';out.mkdir(parents=True,exist_ok=True)
    name=_safe_name(Path(filename).stem if filename else title)+'.'+fmt;target=_resolve(workdir,f'outputs/{name}')
    lines=[line.rstrip() for line in content.splitlines()]
    if fmt=='md':target.write_text(f'# {title}\n\n{content.strip()}\n',encoding='utf-8')
    elif fmt=='docx':
        from docx import Document
        from docx.shared import Pt
        doc=Document();doc.add_heading(title,0)
        for line in lines:
            stripped=line.strip()
            if not stripped:doc.add_paragraph();continue
            level=len(stripped)-len(stripped.lstrip('#'))
            if level:doc.add_heading(stripped[level:].strip(),level=min(3,level))
            elif stripped.startswith(('- ','* ')):doc.add_paragraph(stripped[2:],style='List Bullet')
            else:doc.add_paragraph(stripped)
        for style in doc.styles:
            if hasattr(style,'font'):style.font.name='LXGW WenKai';style.font.size=style.font.size or Pt(11)
        doc.save(target)
    elif fmt=='xlsx':
        from openpyxl import Workbook
        from openpyxl.styles import Font,PatternFill,Alignment
        wb=Workbook();ws=wb.active;ws.title=_safe_name(title)[:31]
        rows=[]
        for line in lines:
            if not line.strip():continue
            cells=[x.strip() for x in line.strip().strip('|').split('|')] if '|' in line else [line.strip()]
            if cells and all(set(x)<=set('-: ') for x in cells):continue
            rows.append(cells)
        for row in rows:ws.append(row)
        if rows:
            for cell in ws[1]:cell.font=Font(name='LXGW WenKai',bold=True,color='FFFFFF');cell.fill=PatternFill('solid',fgColor='365B47');cell.alignment=Alignment(horizontal='center')
            ws.freeze_panes='A2';ws.auto_filter.ref=ws.dimensions
            for column in ws.columns:
                letter=column[0].column_letter;ws.column_dimensions[letter].width=min(42,max(10,max(len(str(c.value or '')) for c in column)+2))
        wb.save(target)
    else:
        from pptx import Presentation
        from pptx.util import Inches,Pt
        deck=Presentation();deck.slide_width=Inches(13.333);deck.slide_height=Inches(7.5)
        cover=deck.slides.add_slide(deck.slide_layouts[0]);cover.shapes.title.text=title;cover.placeholders[1].text='数字人生 · 智能协作生成'
        sections=[];current=('内容提要',[])
        for line in lines:
            if line.startswith('## '):sections.append(current);current=(line[3:].strip(),[])
            elif line.strip():current[1].append(line.lstrip('#- * ').strip())
        sections.append(current)
        for heading,items in sections:
            if not items and heading=='内容提要':continue
            slide=deck.slides.add_slide(deck.slide_layouts[1]);slide.shapes.title.text=heading;frame=slide.placeholders[1].text_frame;frame.clear()
            for i,item in enumerate(items[:8]):p=frame.paragraphs[0] if i==0 else frame.add_paragraph();p.text=item;p.font.size=Pt(22)
        deck.save(target)
    return {'success':True,'format':fmt,'filename':name,'path':f'outputs/{name}','message':f'已生成 {name}'}


# ---------------------------------------------------------------------------
# 卷宗（问策 Run 最终报告）读取 / 列举 / 回填
#
# 卷宗对应一次 GtTaskRun：其正文即该 Run 的最终综合结论（run.final_answer），
# 落盘副本位于团队 outputs 目录，供 V2 / 旧版前端下载。以下函数供 controller
# 直接调用，controller 只负责权限校验与 HTTP 编排，不再自行拼装文件路径。
# ---------------------------------------------------------------------------

def _dossier_report_relpath(run: GtTaskRun) -> str:
    """卷宗落盘文件的团队内相对路径。

    优先复用 run.metadata 中已记录的 final_report_path（与 runService 写入保持
    一致），历史 Run 缺失该字段时退回与 runService 完全相同的确定性文件名，
    确保回填写入与既有写入指向同一文件、互相幂等。
    """
    recorded = (run.metadata or {}).get('final_report_path')
    if isinstance(recorded, str) and recorded.strip():
        return recorded.strip()
    return f'outputs/综合分析报告-run-{run.id}.md'


def _write_report_file(workdir: Path, relpath: str, content: str) -> None:
    """在团队沙箱内原子写入卷宗文件（确定性路径 + os.replace）。"""
    target = _resolve(workdir, relpath)
    target.parent.mkdir(parents=True, exist_ok=True)
    fileUtil.assert_path_within_sandbox(str(target), str(workdir))
    temp_target = target.with_name(f'.{target.name}.{os.getpid()}.tmp')
    temp_target.write_text(content.rstrip() + '\n', encoding='utf-8')
    os.replace(temp_target, target)


async def ensure_dossier_artifact(run: GtTaskRun) -> str | None:
    """确保卷宗落盘文件存在，返回其团队内相对路径；无正文可写时返回 None。

    幂等：文件已存在则直接返回其相对路径；否则用 run.final_answer 回填写入，
    并尽力把相对路径回写进 run.metadata['final_report_path']（best-effort，失败
    不影响本次读取）。这是“历史卷宗无法查看”的核心修复点——H1 之前完成的 Run
    只在 DB 里留有 final_answer 而没有文件，读取时在此按需重建。
    """
    relpath = _dossier_report_relpath(run)
    workdir = await _team_workdir_by_id(run.team_id)
    target = _resolve(workdir, relpath)
    if target.is_file():
        return relpath
    content = (run.final_answer or '').strip()
    if not content:
        return None
    _write_report_file(workdir, relpath, content)
    logger.info('Backfilled dossier artifact: run_id=%s path=%s', run.id, relpath)
    if (run.metadata or {}).get('final_report_path') != relpath:
        try:
            metadata = dict(run.metadata or {})
            metadata['final_report_path'] = relpath
            await gtTaskRunManager.update_run(run.id, metadata=metadata)
        except Exception:
            logger.exception('Failed to persist final_report_path: run_id=%s', run.id)
    return relpath


async def get_dossier(run_id: int) -> dict | None:
    """读取单个卷宗（问策 Run 最终报告）。

    Run 不存在返回 None（controller 据此回 404）。否则返回可查看的卷宗数据：
    - run: 复用 GtTaskRun.to_json() 自动序列化，携带标题、问题、状态、时间、博客
      发布状态等；
    - content: 卷宗正文（最终综合结论 Markdown），历史卷宗同样可读；
    - report_path: 落盘文件的团队内相对路径，供 /files/download.json 下载；文件
      缺失且无正文可回填时为 None；
    - report_ready: 落盘文件当前是否存在（下载链接是否可用）；
    - has_conclusion: 是否已产出结论（final_answer 非空）。
    """
    run = await gtTaskRunManager.get_run(run_id)
    if run is None:
        return None
    report_path = await ensure_dossier_artifact(run)
    content = (run.final_answer or '').strip()
    if report_path and not content:
        # 极少数历史数据：文件在、DB final_answer 为空，则以文件内容为准。
        try:
            content = _resolve(await _team_workdir_by_id(run.team_id), report_path).read_text(
                encoding='utf-8', errors='replace').strip()
        except OSError:
            logger.exception('Failed to read dossier file: run_id=%s path=%s', run.id, report_path)
    return {
        'run': run.to_json(),
        'content': content,
        'report_path': report_path,
        'report_ready': bool(report_path),
        'has_conclusion': bool((run.final_answer or '').strip()),
    }


async def list_dossiers(team_id: int, *, limit: int = 50, owner_user_id: int | None = None) -> list[dict]:
    """列举某团队的卷宗（按 Run 倒序）。

    只做只读汇总，不写盘（避免列举页触发大量回填）。每条给出该卷宗是否可查看：
    - report_ready: 落盘文件当前是否存在；
    - has_conclusion: 是否已产出结论（未落盘但有结论时，get_dossier 会按需重建）。
    """
    runs = await gtTaskRunManager.list_runs(team_id, limit=limit, owner_user_id=owner_user_id)
    try:
        workdir = await _team_workdir_by_id(team_id)
    except ValueError:
        workdir = None
    entries: list[dict] = []
    for run in runs:
        has_conclusion = bool((run.final_answer or '').strip())
        relpath = _dossier_report_relpath(run)
        report_ready = False
        if workdir is not None:
            try:
                report_ready = _resolve(workdir, relpath).is_file()
            except Exception:
                report_ready = False
        entries.append({
            'run': run.to_json(),
            'report_path': relpath if (report_ready or has_conclusion) else None,
            'report_ready': report_ready,
            'has_conclusion': has_conclusion,
        })
    return entries
